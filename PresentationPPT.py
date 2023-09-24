from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Add slides with content
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
slide_1.shapes.title.text = "Slide 1"
slide_1.shapes.placeholders[1].text = "Content for Slide 1"

slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
slide_2.shapes.title.text = "Slide 2"
slide_2.shapes.placeholders[1].text = "Content for Slide 2"

# Save the presentation
presentation.save("new_presentation.pptx")
